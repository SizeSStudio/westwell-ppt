#!/usr/bin/env python3
"""
Westwell PPT Builder  v2.0
──────────────────────────
WestwellPPT class with slide-type methods that match the actual visual
language of Westwell solution decks (airport, factory, seaport references).

Design philosophy (v2):
  • Titles: top-left, left-aligned, navy, ~22pt bold — never centered
  • Thin teal horizontal rule separates title from content
  • Large text / large image / large whitespace — "presentation style"
  • Dark slides for: cover, chapter, statement moments
  • Layout placeholders are always suppressed; content added as free textboxes

Usage:
    import sys; sys.path.insert(0, '~/.claude/skills/westwell-ppt')
    from scripts.pptx_builder import WestwellPPT, preview_pptx

    ppt = WestwellPPT(template='…/PPTTemplate.potx', output='out.pptx')
    ppt.cover("Title", "Subtitle", "2026年4月")
    ppt.agenda(["Chapter 1", "Chapter 2", "Chapter 3"])
    ppt.chapter("01", "Chapter Title", "Supporting line")
    ppt.statement("核心论点：一句话定义这一章的价值")
    ppt.bullets("结论性标题", ["要点A", "要点B", "要点C"])
    ppt.text_image("结论性标题", "Body text…", "/path/image.png")
    ppt.stats("结论性标题", [("95%","指标","说明"), ("<1s","指标","说明")])
    ppt.two_col("标题", "左栏头", "左栏内容", "右栏头", "右栏内容")
    ppt.table("标题", ["Col1","Col2"], [["r1c1","r1c2"]])
    ppt.image("标题", "/path/image.png", "caption")
    ppt.end("期待与您深入探讨")
    ppt.save()
"""

import os, subprocess, tempfile, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colours ───────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x27, 0x43, 0xC6)   # #2743C6  dark blue background
C_NAVY   = RGBColor(0x00, 0x1B, 0xA6)   # #001BA6  title / heading / primary text on light
C_TEAL   = RGBColor(0x00, 0xCA, 0xD4)   # #00CAD4  accent / rule / numbers (non-text)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF0, 0xF4, 0xFF)   # very light blue-grey, card bg
C_MGRAY  = RGBColor(0xD0, 0xD8, 0xF0)   # medium gray, card borders
C_GRAY   = RGBColor(0x55, 0x55, 0x66)   # secondary body text
C_BLACK  = RGBColor(0x1A, 0x1A, 0x2E)   # near-black body text
C_ALTROW = RGBColor(0xE4, 0xEC, 0xFB)   # table alternate row

# ── Geometry (inches) ─────────────────────────────────────────────────────────
SW, SH = 13.333, 7.500      # slide canvas

# Standard content slide title
TL, TT   = 0.906, 0.48      # title left, top
TW, TH   = 11.521, 1.05     # title width (full-bleed reference), height
# Safe title textbox width: caps title so it does not collide with the
# "Make a WELL Change" brand label in the top-right of every content slide.
# Leave ≈2.5" of right margin for the brand zone.
TITLE_SAFE_W     = 9.00
TITLE_MAX_CHARS  = 25       # soft limit per line — warning, not an error

# Teal rule below title (sits just below the visible glyph line at 26pt).
# Previous value 1.58 left a visible gap between title text and rule;
# 1.15 tucks the rule right under the text baseline for a tighter title block.
RULE_Y   = 1.15              # y-position of horizontal rule
RULE_H   = 0.05              # rule height (thin line)

# Content area (below rule, above WMF decoration)
CL       = 0.906
CT       = 1.80              # content top (below rule, with breathing room)
CR       = 12.427
CB       = 5.70              # conservative bottom (WMF circles at 5.933")
CW       = CR - CL           # 11.521"
CH       = CB - CT           # 3.90"

# Editorial framing (eyebrow / subtitle / footnote) — strategy-memo style.
EYEBROW_Y   = 0.18
EYEBROW_H   = 0.28
SUBTITLE_Y  = 1.28
SUBTITLE_H  = 0.62
CT_WITH_SUB = 2.00
FOOTNOTE_Y  = 5.42
FOOTNOTE_H  = 0.30
CB_WITH_FN  = 5.30

# Cover geometry
COV_TL, COV_TT = 0.689, 1.70   # cover title (left half, moved up)
COV_TW, COV_TH = 8.754, 2.705
COV_SUB_Y      = 4.60           # cover subtitle / meta area start

# ── Fonts ─────────────────────────────────────────────────────────────────────
F_TITLE  = 'Encode Sans'
F_BODY   = '思源黑体'
F_ACCENT = 'Encode Sans'


def I(x):
    return Inches(x)


# ── Preview helper ────────────────────────────────────────────────────────────

def preview_pptx(pptx_path: str, out_dir: str, dpi_scale: float = 1.5):
    """
    Convert pptx → PDF (via soffice) → PNG per slide (via pymupdf/fitz).
    Saves files to out_dir as <stem>_slide_NN.png.
    """
    import fitz
    os.makedirs(out_dir, exist_ok=True)
    stem = os.path.splitext(os.path.basename(pptx_path))[0]

    # Convert to PDF
    pdf_path = os.path.join(out_dir, stem + '_preview.pdf')
    subprocess.run(
        ['soffice', '--headless', '--convert-to', 'pdf',
         '--outdir', out_dir, pptx_path],
        check=True, capture_output=True
    )
    # soffice writes to same dir with same stem
    converted = os.path.join(out_dir, stem + '.pdf')
    if os.path.exists(converted) and converted != pdf_path:
        os.rename(converted, pdf_path)

    # Render pages
    doc = fitz.open(pdf_path)
    mat = fitz.Matrix(dpi_scale * 96 / 72, dpi_scale * 96 / 72)
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat)
        pix.save(os.path.join(out_dir, f'{stem}_slide_{i+1:02d}.png'))


# ── Main class ────────────────────────────────────────────────────────────────

class WestwellPPT:
    """Westwell-branded PPTX builder. v2.0 — Westwell visual language."""

    def __init__(self, template: str, output: str):
        self.output = output
        self._pptx_path = self._ensure_pptx(os.path.expanduser(template))
        self.prs = Presentation(self._pptx_path)
        self._layout_map = {l.name: l for l in self.prs.slide_master.slide_layouts}
        self._clear_slides()

    # ── Init helpers ──────────────────────────────────────────────────────────

    def _ensure_pptx(self, path: str) -> str:
        if path.lower().endswith('.pptx'):
            return path
        tmp = tempfile.mkdtemp()
        subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pptx',
             '--outdir', tmp, path],
            check=True, capture_output=True
        )
        base = os.path.splitext(os.path.basename(path))[0]
        out = os.path.join(tmp, base + '.pptx')
        if not os.path.exists(out):
            raise FileNotFoundError(f'Conversion failed: {out}')
        return out

    def _clear_slides(self):
        lst = self.prs.slides._sldIdLst
        for el in list(lst):
            rid = el.get(qn('r:id'))
            lst.remove(el)
            try:
                self.prs.part.drop_rel(rid)
            except Exception:
                pass

    def _layout(self, name: str):
        if name in self._layout_map:
            return self._layout_map[name]
        return self._layout_map.get(
            'custom slide1-light',
            list(self._layout_map.values())[0]
        )

    # ── Low-level drawing primitives ──────────────────────────────────────────

    def _suppress_placeholders(self, slide):
        """Remove all layout placeholders so none show 'Click to add…'."""
        for ph in list(slide.placeholders):
            sp = ph._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)

    def _textbox(self, slide, l, t, w, h, text,
                 size=18, bold=False, color=None,
                 align=PP_ALIGN.LEFT, font=None, wrap=True,
                 line_spacing=None):
        color = color or C_BLACK
        font  = font  or F_BODY
        tb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        if line_spacing:
            from pptx.util import Pt as _Pt
            from pptx.oxml.ns import qn as _qn
            from lxml import etree as _et
            pPr = p._pPr
            if pPr is None:
                pPr = _et.SubElement(p._p, _qn('a:pPr'))
            lnSpc = _et.SubElement(pPr, _qn('a:lnSpc'))
            spcPts = _et.SubElement(lnSpc, _qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing * 100)))
        run = p.add_run()
        run.text = text
        run.font.name  = font
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        return tb

    def _rich_textbox(self, slide, l, t, w, h, text,
                      size=18, color=None, font=None,
                      align=PP_ALIGN.LEFT, wrap=True):
        """
        Textbox with inline **bold** markdown support + `\\n` as paragraph break.
        Any text wrapped in **...** is rendered bold; the rest is regular weight.
        Newlines split text into paragraphs.
        Falls back to _textbox if no markers AND no newlines are found.
        """
        import re
        color = color or C_BLACK
        font  = font  or F_BODY
        has_bold = '**' in text
        has_nl   = '\n' in text
        if not has_bold and not has_nl:
            # No bold markers and no newlines — use simple textbox
            return self._textbox(slide, l, t, w, h, text,
                                 size=size, color=color, font=font,
                                 align=align, wrap=wrap)
        tb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        lines = text.split('\n')
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = align
            parts = re.split(r'\*\*(.*?)\*\*', line)
            if len(parts) == 1 and not parts[0]:
                # Empty line — add a blank run to preserve spacing
                run = p.add_run()
                run.text = ''
                run.font.name  = font
                run.font.size  = Pt(size)
                run.font.color.rgb = color
                continue
            for i, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                run.font.name  = font
                run.font.size  = Pt(size)
                run.font.bold  = (i % 2 == 1)   # odd index = between ** **
                run.font.color.rgb = color
        return tb

    def _para_textbox(self, slide, l, t, w, h, paras, wrap=True):
        """Multi-paragraph textbox. paras = list of dicts."""
        tb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = para.get('align', PP_ALIGN.LEFT)
            if 'space_before' in para:
                p.space_before = Pt(para['space_before'])
            if 'space_after' in para:
                p.space_after = Pt(para['space_after'])
            run = p.add_run()
            run.text = para.get('text', '')
            run.font.name  = para.get('font', F_BODY)
            if 'size'  in para: run.font.size  = Pt(para['size'])
            if 'bold'  in para: run.font.bold  = para['bold']
            if 'color' in para: run.font.color.rgb = para['color']
        return tb

    def _rect(self, slide, l, t, w, h, fill=None, line=None, line_w=0,
              radius=0):
        """Filled rectangle. radius > 0 gives rounded corners (approx via shape)."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if radius > 0:
            # Use RoundedRectangle (shape type 5)
            shape = slide.shapes.add_shape(5, I(l), I(t), I(w), I(h))
            # Set corner radius via XML adj
            sp_el = shape._element
            spPr = sp_el.find(qn('p:spPr'))
            if spPr is not None:
                prstGeom = spPr.find(qn('a:prstGeom'))
                if prstGeom is not None:
                    avLst = prstGeom.find(qn('a:avLst'))
                    if avLst is None:
                        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
                    gd = etree.SubElement(avLst, qn('a:gd'))
                    # radius as fraction of shorter dimension, in 1/100000 units
                    r_val = int(min(radius / min(w, h), 0.5) * 100000)
                    gd.set('name', 'adj')
                    gd.set('fmla', f'val {r_val}')
        else:
            shape = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))

        if fill is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        if line is not None:
            shape.line.color.rgb = line
            if line_w:
                shape.line.width = Pt(line_w)
        else:
            shape.line.fill.background()
        return shape

    def _img_size(self, img_path: str, max_w: float, max_h: float):
        """
        Return (w, h) in inches that fit within max_w × max_h, preserving aspect ratio.
        Reads PNG dimensions from the file header (no PIL dependency).
        Falls back to (max_w, max_h) if the file can't be parsed.
        """
        import struct
        try:
            with open(img_path, 'rb') as f:
                header = f.read(24)
            ext = os.path.splitext(img_path)[1].lower()
            if ext == '.png' and header[:8] == b'\x89PNG\r\n\x1a\n':
                iw = struct.unpack('>I', header[16:20])[0]
                ih = struct.unpack('>I', header[20:24])[0]
            elif ext in ('.jpg', '.jpeg'):
                # JPEG: scan for SOF marker
                with open(img_path, 'rb') as f:
                    data = f.read(65536)
                pos = 2
                iw = ih = 0
                while pos < len(data) - 8:
                    marker = data[pos:pos+2]
                    if marker[0] != 0xFF:
                        break
                    length = struct.unpack('>H', data[pos+2:pos+4])[0]
                    if marker[1] in (0xC0, 0xC1, 0xC2):
                        ih = struct.unpack('>H', data[pos+5:pos+7])[0]
                        iw = struct.unpack('>H', data[pos+7:pos+9])[0]
                        break
                    pos += 2 + length
                if not iw or not ih:
                    return max_w, max_h
            else:
                return max_w, max_h
            scale = min(max_w / iw, max_h / ih)
            return iw * scale, ih * scale
        except Exception:
            return max_w, max_h

    def _hline(self, slide, l, t, w, color=C_TEAL, weight=0.04):
        """Thin horizontal rule (teal by default)."""
        self._rect(slide, l, t, w, weight, fill=color)

    def _title(self, slide, text, dark=False, size=26, eyebrow=''):
        """Standard content-slide title: left-aligned, top-left, +rule below.

        Title width is capped at TITLE_SAFE_W so long titles do not collide
        with the top-right 'Make a WELL Change' brand label. Titles that
        exceed TITLE_MAX_CHARS per line emit a warning — they'll still
        render, but the author should shorten or split them.

        eyebrow (optional): short mono-caps label rendered above the title
        in teal (e.g. 'SECTION 03 · CONTEXT'). Strategy-memo style.
        """
        # Soft length check (first line only — users can split via \n).
        import sys as _sys
        for line in text.split('\n'):
            if len(line) > TITLE_MAX_CHARS:
                print(
                    f'[westwell-ppt] WARN: title line has {len(line)} chars '
                    f'(soft max {TITLE_MAX_CHARS}):\n'
                    f'                     {line!r}\n'
                    f'                     Consider shortening or using \\n '
                    f'to wrap earlier.',
                    file=_sys.stderr,
                )
                break
        # Eyebrow — small bold mono caps in teal above the title
        if eyebrow:
            self._textbox(slide, TL, EYEBROW_Y, TITLE_SAFE_W, EYEBROW_H,
                          eyebrow.upper(), size=11, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
        color = C_WHITE if dark else C_NAVY
        self._textbox(slide, TL, TT, TITLE_SAFE_W, TH,
                      text, size=size, bold=True, color=color,
                      font=F_TITLE, align=PP_ALIGN.LEFT)
        rule_color = C_TEAL if not dark else C_TEAL
        # Short accent rule (left-anchored, ~1.2" wide) — Westwell signature
        self._rect(slide, TL, RULE_Y, 1.20, RULE_H, fill=rule_color)
        # Soft long rule under title, light gray
        self._rect(slide, TL + 1.20, RULE_Y + 0.012,
                   TW - 1.20, 0.015,
                   fill=C_MGRAY if not dark else RGBColor(0x3A, 0x56, 0xD0))

    # ── Editorial framing helpers ─────────────────────────────────────────────

    def _subtitle_block(self, slide, text, dark=False):
        """14pt gray subtitle paragraph below title rule (y=1.28, h=0.62).
        Renders via _rich_textbox (supports **bold** + \\n)."""
        if not text:
            return None
        color = C_LGRAY if dark else C_GRAY
        return self._rich_textbox(
            slide, TL, SUBTITLE_Y, TITLE_SAFE_W, SUBTITLE_H,
            text, size=14, color=color, font=F_BODY,
        )

    def _footnote_strip(self, slide, text, dark=False):
        """Italic 12pt gray 'so-what' strip at y=5.42, h=0.30.
        Use python-pptx low-level to set italic=True."""
        if not text:
            return None
        color = C_LGRAY if dark else C_GRAY
        tb = slide.shapes.add_textbox(
            I(TL), I(FOOTNOTE_Y), I(TITLE_SAFE_W), I(FOOTNOTE_H))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = F_BODY
        run.font.size = Pt(12)
        run.font.italic = True
        run.font.color.rgb = color
        return tb

    def _content_bounds(self, subtitle='', footnote='', bottom_callout=None):
        """Returns (ct, cb, ch) for dynamic content area.
        ct = CT_WITH_SUB if subtitle else CT
        cb = 4.70 if bottom_callout else CB_WITH_FN if footnote else CB
        Returns (ct, cb, cb - ct)."""
        ct = CT_WITH_SUB if subtitle else CT
        if bottom_callout:
            cb = 4.70
        elif footnote:
            cb = CB_WITH_FN
        else:
            cb = CB
        return ct, cb, cb - ct

    def _bottom_callout(self, slide, callout, dark=False):
        """Render callout dict={'label': 'BOTTOM LINE', 'text': '...', 'dark': True}
        at y=4.80, h=0.85. Dark variant: C_NAVY bg + C_TEAL label + C_WHITE text.
        Light variant: C_LGRAY bg + navy left-bar + C_TEAL label + C_NAVY text.
        Label rendered at 11pt mono caps. Main text 17pt rich_textbox."""
        if not callout:
            return None
        y = 4.80
        h = 0.85
        w = TW
        x = TL
        label = callout.get('label', 'BOTTOM LINE')
        text  = callout.get('text', '')
        is_dark = callout.get('dark', dark)
        if is_dark:
            # Dark variant
            self._rect(slide, x, y, w, h, fill=C_NAVY, radius=0.06)
            # Optional thin teal left bar
            self._rect(slide, x, y, 0.08, h, fill=C_TEAL)
            label_color = C_TEAL
            text_color  = C_WHITE
        else:
            # Light variant
            self._rect(slide, x, y, w, h, fill=C_LGRAY, radius=0.06)
            self._rect(slide, x, y, 0.08, h, fill=C_NAVY)
            label_color = C_TEAL
            text_color  = C_NAVY
        # Label — mono caps 11pt teal
        self._textbox(slide, x + 0.25, y + 0.10, w - 0.40, 0.28,
                      label.upper(), size=11, bold=True, color=label_color,
                      font=F_ACCENT, align=PP_ALIGN.LEFT)
        # Main text — 17pt rich textbox
        self._rich_textbox(slide, x + 0.25, y + 0.38, w - 0.40, h - 0.44,
                           text, size=17, color=text_color, font=F_BODY)
        return None

    def _attach_notes(self, slide, notes=''):
        """If notes non-empty, set slide.notes_slide.notes_text_frame.text = notes."""
        if not notes:
            return
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass

    def _new_slide(self, layout_name: str, dark=False,
                   density: str = 'compact') -> object:
        """Add slide, suppress all placeholders.

        density controls which content master is used:
          • 'compact'  → `custom slide2-{mode}` (slim bottom decoration,
            the Westwell-preferred default — lets content breathe without
            the heavy decorative footer competing visually)
          • 'standard' → `custom slide1-{mode}` (full decorative bottom;
            use sparingly when you want extra visual rhythm on an
            otherwise sparse page)

        Default is 'compact' because the full decorative bottom tends to
        out-weigh the content on most pages. Only pass density='standard'
        when you explicitly want the heavier footer for visual variety.
        """
        if density == 'standard':
            name = 'custom slide1-dark' if dark else layout_name
        else:
            name = 'custom slide2-dark' if dark else 'custom slide2-light'
        slide = self.prs.slides.add_slide(self._layout(name))
        self._suppress_placeholders(slide)
        return slide

    # ── Public slide methods ───────────────────────────────────────────────────

    def cover(self, title: str, subtitle: str = '',
              context: str = '', date: str = ''):
        """
        Cover slide. Dark blue background, Westwell geometric art right side.
        Title occupies left ~65% of the slide.
        All layout placeholders are suppressed — content added as free textboxes.
        """
        slide = self.prs.slides.add_slide(self._layout('标题幻灯片'))
        self._suppress_placeholders(slide)

        # Main title (left half, upper-center)
        self._textbox(slide, COV_TL, COV_TT, COV_TW, COV_TH,
                      title, size=40, bold=True, color=C_WHITE,
                      font=F_TITLE, align=PP_ALIGN.LEFT)

        # Thin teal rule below title
        lines = title.count('\n') + 1
        rule_y = COV_TT + lines * 0.68 + 0.15
        self._hline(slide, COV_TL, rule_y, 5.5, C_TEAL, 0.05)

        # Subtitle / context / date stacked below rule
        y = COV_SUB_Y
        if subtitle:
            self._textbox(slide, COV_TL, y, COV_TW, 0.55,
                          subtitle, size=18, color=C_LGRAY,
                          font=F_BODY, align=PP_ALIGN.LEFT)
            y += 0.58
        if context:
            self._textbox(slide, COV_TL, y, COV_TW, 0.40,
                          context, size=14, color=C_TEAL,
                          font=F_TITLE, align=PP_ALIGN.LEFT)
            y += 0.43
        if date:
            self._textbox(slide, COV_TL, y, COV_TW, 0.35,
                          date, size=13, color=C_LGRAY,
                          font=F_TITLE, align=PP_ALIGN.LEFT)
        return slide

    def agenda(self, items: list):
        """
        Table-of-contents slide.
        items: list of strings OR (num, title) OR (num, title, subtitle) tuples.
        Layout: 2-row card grid on white background.
        """
        slide = self._new_slide('custom slide1-light')

        # "目录" header
        self._textbox(slide, TL, TT, 3.0, TH,
                      '目 录', size=26, bold=True, color=C_NAVY,
                      font=F_TITLE, align=PP_ALIGN.LEFT)
        self._hline(slide, TL, RULE_Y, TW, C_TEAL, 0.04)

        n = len(items)
        cols = min(3, n)
        rows = math.ceil(n / cols)

        GAP  = 0.20
        card_w = (CW - GAP * (cols - 1)) / cols
        card_h = (CH - GAP * (rows - 1)) / rows

        for i, item in enumerate(items):
            if isinstance(item, str):
                num = f'{i+1:02d}'
                ttl = item
                sub = ''
            elif len(item) >= 3:
                num, ttl, sub = str(item[0]), str(item[1]), str(item[2])
            else:
                num = str(item[0]) if len(item) > 0 else f'{i+1:02d}'
                ttl = str(item[1]) if len(item) > 1 else ''
                sub = ''

            col_i = i % cols
            row_i = i // cols
            x = CL + col_i * (card_w + GAP)
            y = CT + row_i * (card_h + GAP)

            # Card background (rounded rect)
            self._rect(slide, x, y, card_w, card_h,
                       fill=C_LGRAY, radius=0.12)
            # Left teal accent strip
            self._rect(slide, x, y, 0.06, card_h, fill=C_TEAL)
            # Chapter number
            self._textbox(slide, x + 0.18, y + 0.18, card_w - 0.28, 0.70,
                          num, size=44, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # Chapter title
            self._textbox(slide, x + 0.18, y + 0.80, card_w - 0.28,
                          card_h - 0.90,
                          ttl, size=16, bold=True, color=C_NAVY,
                          font=F_BODY, align=PP_ALIGN.LEFT)
            if sub:
                self._textbox(slide, x + 0.18, y + card_h - 0.45,
                              card_w - 0.28, 0.40,
                              sub, size=11, color=C_GRAY,
                              font=F_BODY, align=PP_ALIGN.LEFT)
        return slide

    def chapter(self, num: str, title: str,
                subtitle: str = '', highlight: str = ''):
        """
        Chapter separator. Uses 'agenda slide 2' layout — dark blue + colorful
        right-side geometric art inherited from template.
        Content occupies the left ~55% of the slide so the art stays visible.
        """
        slide = self.prs.slides.add_slide(self._layout('agenda slide 2'))
        self._suppress_placeholders(slide)

        # Constrain content to left side so template's right-side art stays visible.
        # Art extends into bottom-left too, so keep text above y~4.6.
        cx_l = 0.906
        cx_w = 5.30

        # Large chapter number — white on dark bg (text-color rule:
        # navy / black / white only).
        self._textbox(slide, cx_l, 0.90, 3.0, 1.50,
                      num, size=100, bold=True, color=C_WHITE,
                      font=F_ACCENT, align=PP_ALIGN.LEFT)

        # Thin teal rule
        self._hline(slide, cx_l, 2.55, 4.20, C_TEAL, 0.06)

        # Chapter title — white, below number
        lines = title.split('\n')
        self._para_textbox(
            slide, cx_l, 2.75, cx_w, 1.80,
            [{'text': ln, 'size': 30, 'bold': True,
              'color': C_WHITE, 'font': F_TITLE,
              'space_after': 4}
             for ln in lines]
        )
        # subtitle/highlight intentionally ignored: template art occupies
        # the lower-left region and any text there becomes unreadable.
        return slide

    def statement(self, text: str, label: str = '', dark: bool = True,
                  density: str = 'compact',
                  eyebrow: str = '', subtitle: str = '',
                  footnote: str = '', notes: str = ''):
        """
        Full-slide statement. A single bold claim, centered.
        Use for pivotal moments: the core thesis, a key data insight,
        or the opening of a chapter.
        label: small caption in top-left (e.g. section context)

        Editorial framing (eyebrow/subtitle/footnote/notes) is supported
        but adds weight to what should be a minimal canvas — use sparingly.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        if eyebrow:
            # Eyebrow at top-left on a statement slide too
            self._textbox(slide, TL, EYEBROW_Y, TITLE_SAFE_W, EYEBROW_H,
                          eyebrow.upper(), size=11, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
        if subtitle:
            self._subtitle_block(slide, subtitle, dark)
        body_color = C_WHITE if dark else C_NAVY

        if label:
            # Centered caption above the main statement, with teal brackets
            self._textbox(slide, TL, 1.80, TW, 0.45,
                          f'— {label} —', size=15, color=C_TEAL,
                          font=F_TITLE, align=PP_ALIGN.CENTER)

        # Each line as its own centered paragraph — avoids alignment drift
        lines = [l for l in text.split('\n') if l.strip()]
        stmt_t = 2.70 if label else 2.40
        self._para_textbox(
            slide, TL, stmt_t, TW, 3.20,
            [{'text': ln, 'size': 36, 'bold': True,
              'color': body_color, 'font': F_BODY,
              'align': PP_ALIGN.CENTER, 'space_after': 18}
             for ln in lines]
        )
        # Short teal accent rule below statement
        self._rect(slide, SW / 2 - 0.6, 5.20, 1.2, 0.05, fill=C_TEAL)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def bullets(self, title: str, points: list,
                dark: bool = False, body_size: int = 19,
                density: str = 'compact',
                eyebrow: str = '', subtitle: str = '',
                footnote: str = '', notes: str = ''):
        """
        Bulleted list slide.
        points: list of strings. Each is one bullet.
        Design: left-aligned title + teal rule + clean bullet items.
        density: 'standard' (default) or 'compact' — use compact for 5+
        bullets or when the full bottom decoration would crowd content.

        Editorial framing:
          • eyebrow  — short mono-caps label above title (e.g. 'WHY · 4 POINTS')
          • subtitle — 14pt gray paragraph below title rule
          • footnote — italic 12pt 'so-what' strip at the bottom
          • notes    — speaker notes
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        body_color = C_WHITE if dark else C_BLACK
        n = len(points)
        gap = 0.10
        item_h = min(0.72, (ch - gap * (n - 1)) / n)

        for i, pt in enumerate(points):
            y = ct + i * (item_h + gap)
            # Teal square marker — vertically centred in item height
            marker_size = 0.10
            marker_y = y + (item_h - marker_size) / 2
            self._rect(slide, CL, marker_y, marker_size, marker_size,
                       fill=C_TEAL)
            # Bullet text — supports **bold** inline markers
            self._rich_textbox(slide, CL + 0.22, y, CW - 0.24, item_h,
                               pt, size=body_size, color=body_color,
                               font=F_BODY)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def text_image(self, title: str, body: str, img_path: str = None,
                   dark: bool = False, body_size: int = 16,
                   density: str = 'compact',
                   eyebrow: str = '', subtitle: str = '',
                   footnote: str = '', notes: str = ''):
        """
        Left text (32%) + right image (64%), side by side.

        Extended content area (CB_IMG=6.60) gives both columns 4.80" of
        vertical room — ~23% taller than the standard CB. The image
        column is deliberately wide so images (architecture diagrams,
        flowcharts) are rendered as large as their aspect ratio allows.

        img_path=None → styled placeholder box.
        Supports **bold** inline markers in body text.
        Image is aspect-ratio-aware (never distorted).

        Editorial framing (eyebrow/subtitle/footnote/notes) is supported
        but does NOT auto-shrink the image/body area — when using
        subtitle or footnote, author should trim body content accordingly
        to avoid overlap with the footnote strip.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        body_color = C_WHITE if dark else C_BLACK

        # Extended content area — match image() so side-by-side layouts
        # don't waste the 0.9" below the standard CB.
        CB_IMG = 6.60
        content_h = CB_IMG - CT            # 4.80

        text_w = CW * 0.32                 # 3.69" (narrower than before)
        gap    = CW * 0.03                 # 0.35"
        img_l  = CL + text_w + gap         # 5.00
        img_w  = CW - text_w - gap         # 7.48" (wider than before)

        # Body text — supports **bold** markers
        self._rich_textbox(slide, CL, CT, text_w, content_h,
                           body, size=body_size, color=body_color,
                           font=F_BODY)

        # Light card backdrop behind image (only on light slides, image present)
        if img_path and os.path.exists(img_path) and not dark:
            self._rect(slide, img_l, CT, img_w, content_h,
                       fill=C_LGRAY, radius=0.10)

        # Image (right) — aspect-ratio-aware with inner padding
        if img_path and os.path.exists(img_path):
            pad = 0.22
            w, h = self._img_size(img_path, img_w - pad * 2, content_h - pad * 2)
            x_off = (img_w - w) / 2
            y_off = (content_h - h) / 2
            slide.shapes.add_picture(
                img_path, I(img_l + x_off), I(CT + y_off), I(w), I(h))
        else:
            # Styled placeholder
            ph_fill = RGBColor(0x1A, 0x2E, 0x7A) if dark else C_LGRAY
            self._rect(slide, img_l, CT, img_w, content_h,
                       fill=ph_fill, radius=0.12)
            self._rect(slide, img_l, CT, img_w, 0.06, fill=C_TEAL)
            lbl_color = C_TEAL
            self._textbox(slide, img_l, CT + content_h / 2 - 0.35, img_w, 0.50,
                          '[ 示意图 ]', size=15, bold=False,
                          color=lbl_color, font=F_TITLE,
                          align=PP_ALIGN.CENTER)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def stats(self, title: str, stats: list, dark: bool = False,
              density: str = 'compact',
              eyebrow: str = '', subtitle: str = '',
              footnote: str = '', notes: str = ''):
        """
        KPI / metric slide.
        stats: list of (value, label, description) tuples. 2–4 items.
        Large teal value + navy label + gray description.

        Value font size auto-shrinks so the longest value fits cleanly:
        mixed values like ("55万", "15–20 分钟", "数十架") all render at
        the same reduced size instead of the long one shrinking alone.

        density: 'standard' or 'compact' — use compact for info-heavy KPI
        pages where the bottom Westwell decoration would compete.

        Editorial framing (eyebrow/subtitle/footnote/notes) is supported
        but does NOT auto-shrink the KPI card area — consider reducing
        description length when using subtitle or footnote.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        # Auto-shrink value font size so all values render at the same
        # scale (longest one dictates). Chinese and ASCII mixed values are
        # treated by character count — conservative but visually reliable.
        max_len = max(len(str(v)) for v, _, _ in stats)
        if   max_len <= 4:  val_size = 54
        elif max_len <= 6:  val_size = 42
        elif max_len <= 8:  val_size = 34
        elif max_len <= 10: val_size = 28
        else:               val_size = 24

        n = len(stats)
        gap = 0.20
        card_w = (CW - gap * (n - 1)) / n
        card_h = CH

        for i, (val, lbl, desc) in enumerate(stats):
            x = CL + i * (card_w + gap)
            # Light card background on light slides
            if not dark:
                self._rect(slide, x, CT, card_w, card_h,
                           fill=C_LGRAY, radius=0.10)

            # Thin top accent
            self._rect(slide, x, CT, card_w, 0.06, fill=C_TEAL)

            # Value — large KPI number. Text-color rule: navy on light,
            # white on dark (no teal text). wrap=False so values like
            # "15–20 min" stay on a single line.
            val_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, x, CT + 0.25, card_w, 1.30,
                          val, size=val_size, bold=True, color=val_color,
                          font=F_ACCENT, align=PP_ALIGN.CENTER, wrap=False)
            # Label — smaller supporting text, same color rule as value
            lbl_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, x, CT + 1.62, card_w, 0.45,
                          lbl, size=17, bold=True, color=lbl_color,
                          font=F_BODY, align=PP_ALIGN.CENTER)
            # Thin separator
            self._hline(slide, x + 0.55, CT + 2.18, card_w - 1.1,
                        C_TEAL if not dark else C_WHITE, 0.025)
            # Description
            self._textbox(slide, x + 0.20, CT + 2.32, card_w - 0.40, 1.50,
                          desc, size=13, color=C_GRAY if not dark else C_LGRAY,
                          font=F_BODY, align=PP_ALIGN.CENTER, wrap=True)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def two_col(self, title: str,
                left_head: str, left_body: str,
                right_head: str, right_body: str,
                dark: bool = False, density: str = 'compact',
                eyebrow: str = '', subtitle: str = '',
                footnote: str = '', bottom_callout: dict = None,
                notes: str = ''):
        """
        Two-column layout. Heading + body in each column.
        A thin vertical teal rule separates the columns.
        density: 'standard' or 'compact' — use compact when each column
        has 6+ lines of body text.

        Editorial framing kwargs:
          • eyebrow, subtitle, footnote, notes — see bullets().
          • bottom_callout — dict {label, text, dark?} rendered as a
            strip at y=4.80; suppresses footnote when present.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, bottom_callout)

        col_w = (CW - 0.30) / 2   # gap between cols = 0.30"
        r_x   = CL + col_w + 0.30

        for x, head, body in [(CL, left_head, left_body),
                               (r_x, right_head, right_body)]:
            # Card background
            card_color = RGBColor(0x1A, 0x2E, 0x7A) if dark else C_LGRAY
            self._rect(slide, x, ct, col_w, ch, fill=card_color, radius=0.10)
            # Top accent strip (teal — non-text accent is allowed)
            self._rect(slide, x, ct, col_w, 0.06, fill=C_TEAL)
            # Column heading — text-color rule: navy on light, white on dark
            head_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, x + 0.28, ct + 0.22, col_w - 0.40, 0.55,
                          head, size=19, bold=True, color=head_color,
                          font=F_BODY, align=PP_ALIGN.LEFT)
            # Thin rule under heading
            self._hline(slide, x + 0.28, ct + 0.82, col_w - 0.56, C_TEAL, 0.022)
            # Column body — supports **bold** markers
            body_color = C_LGRAY if dark else C_BLACK
            self._rich_textbox(slide, x + 0.28, ct + 1.00, col_w - 0.40,
                               ch - 1.15, body, size=17, color=body_color,
                               font=F_BODY)
        if bottom_callout:
            self._bottom_callout(slide, bottom_callout, dark)
        self._footnote_strip(slide, footnote if not bottom_callout else '', dark)
        self._attach_notes(slide, notes)
        return slide

    def three_col(self, title: str, columns: list, dark: bool = False,
                  density: str = 'compact',
                  eyebrow: str = '', subtitle: str = '',
                  footnote: str = '', bottom_callout: dict = None,
                  notes: str = ''):
        """
        Three-column insight card layout.
        columns: list of 3 dicts {"head": str, "body": str}
        Best for: executive summary pillars, phase overview, 3-way comparison.
        Supports **bold** inline markers in body text.

        Editorial framing kwargs:
          • eyebrow, subtitle, footnote, notes — see bullets().
          • bottom_callout — dict {label, text, dark?}; suppresses footnote.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, bottom_callout)

        GAP    = 0.18
        col_w  = (CW - GAP * 2) / 3

        for i, col in enumerate(columns[:3]):
            x = CL + i * (col_w + GAP)
            card_color = RGBColor(0x1A, 0x2E, 0x7A) if dark else C_LGRAY
            # Card background
            self._rect(slide, x, ct, col_w, ch, fill=card_color, radius=0.10)
            # Teal top accent strip
            self._rect(slide, x, ct, col_w, 0.06, fill=C_TEAL)
            # Column heading — text-color rule: navy on light, white on dark
            head_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, x + 0.24, ct + 0.22, col_w - 0.36, 0.55,
                          col.get('head', ''), size=18, bold=True,
                          color=head_color, font=F_BODY, align=PP_ALIGN.LEFT)
            # Rule under heading
            self._hline(slide, x + 0.24, ct + 0.82, col_w - 0.48, C_TEAL, 0.022)
            # Body text — supports **bold** markers
            body_color = C_LGRAY if dark else C_BLACK
            self._rich_textbox(slide, x + 0.24, ct + 1.00, col_w - 0.36,
                               ch - 1.15, col.get('body', ''),
                               size=16, color=body_color, font=F_BODY)
        if bottom_callout:
            self._bottom_callout(slide, bottom_callout, dark)
        self._footnote_strip(slide, footnote if not bottom_callout else '', dark)
        self._attach_notes(slide, notes)
        return slide

    def table(self, title: str, headers: list, rows: list,
              dark: bool = False, bold_col: int = None,
              density: str = 'compact', body_size: int = 11,
              eyebrow: str = '', subtitle: str = '',
              footnote: str = '', notes: str = ''):
        """
        Data table. Navy header row, alternating body rows.

        body_size: body cell font size (pt). Default 11 is tuned for
        dense tables (6+ cols × 8+ rows). For sparse tables (3–5 cols ×
        4–6 rows) bump to 14–16 so text has visual weight matching its
        generous cell space. Header row scales to body_size + 2.

        density: 'standard' or 'compact' — compact uses a slimmer bottom
        decoration so 6+ row tables do not crowd the footer.

        Editorial framing (eyebrow/subtitle/footnote/notes) is supported
        but does NOT auto-shrink the table area — trim row count when
        using subtitle or footnote to avoid overlap.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        n_rows = 1 + len(rows)
        n_cols = len(headers)
        tbl = slide.shapes.add_table(
            n_rows, n_cols, I(CL), I(CT), I(CW), I(CH - 0.05)
        ).table

        col_w = I(CW) / n_cols
        for col in tbl.columns:
            col.width = int(col_w)

        header_size = body_size + 2

        # Header row
        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = hdr
            p = cell.text_frame.paragraphs[0]
            p.font.name = F_TITLE
            p.font.size = Pt(header_size)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_NAVY

        # Data rows — text-color rule: navy (bold_col) / black (normal),
        # no teal text.
        for ri, row_data in enumerate(rows):
            bg = C_WHITE if ri % 2 == 0 else C_ALTROW
            for ci, val in enumerate(row_data):
                cell = tbl.cell(ri + 1, ci)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.name = F_BODY
                p.font.size = Pt(body_size)
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                if bold_col is not None and ci == bold_col:
                    p.font.bold = True
                    p.font.color.rgb = C_NAVY
                else:
                    p.font.color.rgb = C_BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def image(self, title: str, img_path: str,
              caption: str = '', body: str = '',
              dark: bool = False, density: str = 'compact',
              eyebrow: str = '', subtitle: str = '',
              footnote: str = '', notes: str = ''):
        """
        Image slide with optional narrative body below.

          • caption — short single-line source note (12pt gray, centered).
          • body    — multi-line narrative explaining what the image shows
                      and which concrete scenario it addresses (15pt,
                      left-aligned, supports **bold** markers). Prefer
                      body over caption when the image needs context.

        Layout:
          • Image is TOP-aligned (not vertically centered) so large
            images get maximum height and the text below reads as a
            natural continuation, not a stranded afterthought.
          • Content area extends down to CB_IMG (~6.60"), deeper than
            the standard CB, because compact layout leaves that region
            clean and an image is the whole point of this slide.

        Editorial framing (eyebrow/subtitle/footnote/notes) is supported
        but does NOT auto-shrink the image area — trim image height or
        body when using subtitle or footnote.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        CB_IMG = 6.60             # extended bottom for image slides
        content_h = CB_IMG - CT   # 4.80" vs standard 3.90"

        if body:
            bottom_h = 1.05       # 3 lines at 15pt + breathing room
        elif caption:
            bottom_h = 0.35
        else:
            bottom_h = 0.0
        avail_h = content_h - bottom_h

        if img_path and os.path.exists(img_path):
            w, h = self._img_size(img_path, CW, avail_h)
            # Horizontally center; TOP-align vertically
            x_off = (CW - w) / 2
            slide.shapes.add_picture(
                img_path, I(CL + x_off), I(CT), I(w), I(h))

        if body:
            body_color = C_WHITE if dark else C_BLACK
            self._rich_textbox(
                slide, CL, CB_IMG - bottom_h, CW, bottom_h,
                body, size=15, color=body_color, font=F_BODY,
            )
        elif caption:
            self._textbox(
                slide, CL, CB_IMG - 0.32, CW, 0.32,
                caption, size=12, color=C_GRAY,
                font=F_BODY, align=PP_ALIGN.CENTER,
            )
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def text(self, title: str, body: str, dark: bool = False,
             body_size: int = 20, density: str = 'compact',
             eyebrow: str = '', subtitle: str = '',
             footnote: str = '', notes: str = ''):
        """
        Simple text slide. Title + rule + flowing body text.
        Use for a single statement backed by a prose explanation.

        Editorial framing (eyebrow/subtitle/footnote/notes) is supported
        but does NOT auto-shrink the body area — trim body when using
        subtitle or footnote to avoid overlap.
        """
        slide = self._new_slide(
            'custom slide1-light', dark=dark, density=density)
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        body_color = C_WHITE if dark else C_BLACK
        self._textbox(slide, CL, CT, CW, CH,
                      body, size=body_size, color=body_color,
                      font=F_BODY, wrap=True)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def end(self, text: str = 'THANK YOU', subtitle: str = ''):
        """
        End slide. Uses 'end slide' layout (dark blue + left colorful art).
        Text is centered on the right two-thirds of the slide.
        """
        slide = self.prs.slides.add_slide(self._layout('end slide'))
        self._suppress_placeholders(slide)

        lines = text.split('\n')
        cx = SW / 2
        self._para_textbox(
            slide, cx - 3.5, 2.50, 7.0, 2.0,
            [{'text': ln, 'size': 34, 'bold': True,
              'color': C_WHITE, 'font': F_TITLE,
              'align': PP_ALIGN.CENTER, 'space_after': 8}
             for ln in lines]
        )
        if subtitle:
            self._textbox(slide, cx - 3.5, 4.55, 7.0, 0.60,
                          subtitle, size=16, color=C_LGRAY,
                          font=F_BODY, align=PP_ALIGN.CENTER)
        return slide

    # ── Compositional layouts (strategy-memo patterns) ────────────────────────

    def pyramid(self, title: str, tiers: list, caption: str = '',
                dark: bool = False,
                eyebrow: str = '', subtitle: str = '',
                footnote: str = '', notes: str = ''):
        """Hierarchy pyramid — 3 horizontal bars, widest at bottom.

        tiers: list of 3 dicts ordered BOTTOM→TOP, each with
               {label, en, sub} keys.
        caption (legacy): rendered bottom-center 12pt gray if no footnote.
        """
        assert len(tiers) == 3, 'pyramid requires exactly 3 tiers'
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        widths  = [9.0, 7.0, 5.0]                 # BOTTOM → TOP widths
        shades  = [RGBColor(0x00, 0x13, 0x6C), C_NAVY, C_DARK]
        tier_labels = ['TIER 01', 'TIER 02', 'TIER 03']

        bar_h = 0.90
        gap   = 0.08
        total = bar_h * 3 + gap * 2
        # Vertically center in content area
        start_y = ct + max(0, (ch - total) / 2)

        for i, tier in enumerate(tiers):   # i=0 bottom, 2 top
            # Bottom first: y increases downward so bottom tier is lowest
            y = start_y + (2 - i) * (bar_h + gap)
            w = widths[i]
            x = CL + (CW - w) / 2
            # Bar background
            self._rect(slide, x, y, w, bar_h, fill=shades[i], radius=0.06)
            # Left teal accent strip
            self._rect(slide, x, y, 0.08, bar_h, fill=C_TEAL)
            # TIER label (bottom=TIER 1, top=TIER 3)
            self._textbox(slide, x + 0.22, y + 0.08, w - 0.30, 0.28,
                          tier_labels[i], size=11, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # CN label (18pt bold white)
            self._textbox(slide, x + 0.22, y + 0.32, w - 0.30, 0.35,
                          tier.get('label', ''), size=18, bold=True,
                          color=C_WHITE, font=F_BODY, align=PP_ALIGN.LEFT)
            # sub (12pt light gray)
            if tier.get('sub'):
                self._textbox(slide, x + 0.22, y + 0.62, (w - 0.30) / 2, 0.28,
                              tier.get('sub', ''), size=12, color=C_LGRAY,
                              font=F_BODY, align=PP_ALIGN.LEFT)
            # EN line (12pt light gray, right side)
            if tier.get('en'):
                self._textbox(slide, x + (w / 2), y + 0.62, (w / 2) - 0.22, 0.28,
                              tier.get('en', ''), size=12, color=C_LGRAY,
                              font=F_ACCENT, align=PP_ALIGN.RIGHT)
        # Legacy caption (bottom-center) only when no footnote
        if caption and not footnote:
            self._textbox(slide, CL, 5.10, CW, 0.30,
                          caption, size=12, color=C_GRAY,
                          font=F_BODY, align=PP_ALIGN.CENTER)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def value_chain(self, title: str, steps: list,
                    highlight_last: bool = True, dark: bool = False,
                    eyebrow: str = '', subtitle: str = '',
                    footnote: str = '', notes: str = ''):
        """3–5 equal-width columns with STEP 0N labels; first step navy
        vertical rule, rest teal.

        steps: list of dicts {title, body}."""
        n = len(steps)
        assert 3 <= n <= 5, 'value_chain requires 3-5 steps'
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        GAP = 0.20
        col_w = (CW - GAP * (n - 1)) / n

        for i, step in enumerate(steps):
            x = CL + i * (col_w + GAP)
            if highlight_last and not dark and i == n - 1:
                self._rect(slide, x, ct, col_w, ch, fill=C_LGRAY, radius=0.08)
            # Left vertical rule
            rule_color = C_NAVY if i == 0 else C_TEAL
            self._rect(slide, x, ct, 0.04, ch, fill=rule_color)
            # STEP label
            self._textbox(slide, x + 0.18, ct + 0.05, col_w - 0.24, 0.28,
                          f'STEP {i+1:02d}', size=11, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # Title
            title_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, x + 0.18, ct + 0.32, col_w - 0.24, 0.50,
                          step.get('title', ''), size=17, bold=True,
                          color=title_color, font=F_BODY, align=PP_ALIGN.LEFT)
            # Body
            body_color = C_LGRAY if dark else C_GRAY
            self._rich_textbox(slide, x + 0.18, ct + 0.90, col_w - 0.24,
                               ch - 1.0, step.get('body', ''),
                               size=13, color=body_color, font=F_BODY)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def control_matrix(self, title: str, must_control: list,
                       can_partner: list, principle: str = '',
                       dark: bool = False,
                       eyebrow: str = '', subtitle: str = '',
                       footnote: str = '', notes: str = ''):
        """Two-panel: 57% dark 'MUST CONTROL' left, 43% light 'CAN PARTNER' right.

        must_control / can_partner: lists of strings."""
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        left_w  = CW * 0.57
        right_w = CW - left_w - 0.18
        right_x = CL + left_w + 0.18

        # Left panel — dark navy
        self._rect(slide, CL, ct, left_w, ch, fill=C_NAVY, radius=0.08)
        self._textbox(slide, CL + 0.28, ct + 0.20, left_w - 0.40, 0.28,
                      'MUST CONTROL', size=11, bold=True, color=C_TEAL,
                      font=F_ACCENT, align=PP_ALIGN.LEFT)
        self._textbox(slide, CL + 0.28, ct + 0.48, left_w - 0.40, 0.55,
                      '必须自控的能力', size=19, bold=True, color=C_WHITE,
                      font=F_BODY, align=PP_ALIGN.LEFT)
        self._hline(slide, CL + 0.28, ct + 1.06, left_w - 0.56, C_TEAL, 0.025)
        # Numbered items
        item_y = ct + 1.22
        n_left = len(must_control)
        item_h_left = min(0.68, (ch - 1.30) / max(n_left, 1))
        for i, it in enumerate(must_control):
            y = item_y + i * item_h_left
            self._textbox(slide, CL + 0.28, y, 0.60, item_h_left,
                          f'{i+1:02d}', size=13, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            self._rich_textbox(slide, CL + 0.85, y, left_w - 1.05, item_h_left,
                               it, size=14, color=C_WHITE, font=F_BODY)

        # Right panel — light gray (or medium navy in dark mode)
        right_fill = RGBColor(0x1A, 0x2E, 0x7A) if dark else C_LGRAY
        self._rect(slide, right_x, ct, right_w, ch, fill=right_fill, radius=0.08)
        lbl_color = C_LGRAY if dark else C_GRAY
        self._textbox(slide, right_x + 0.24, ct + 0.20, right_w - 0.36, 0.28,
                      'CAN PARTNER', size=11, bold=True, color=lbl_color,
                      font=F_ACCENT, align=PP_ALIGN.LEFT)
        title_color = C_WHITE if dark else C_NAVY
        self._textbox(slide, right_x + 0.24, ct + 0.48, right_w - 0.36, 0.50,
                      '可以合作获取', size=17, bold=True, color=title_color,
                      font=F_BODY, align=PP_ALIGN.LEFT)
        self._hline(slide, right_x + 0.24, ct + 1.06, right_w - 0.48, C_TEAL, 0.022)
        # Bullet items
        body_color = C_LGRAY if dark else C_BLACK
        # Compute available items region — if principle given, reserve 0.80" at bottom
        principle_h = 0.80 if principle else 0.0
        items_y0 = ct + 1.22
        items_avail = ch - 1.30 - principle_h
        n_right = len(can_partner)
        item_h_right = min(0.56, items_avail / max(n_right, 1))
        for i, it in enumerate(can_partner):
            y = items_y0 + i * item_h_right
            # Teal dot
            dot = 0.12
            self._rect(slide, right_x + 0.28, y + (item_h_right - dot) / 2,
                       dot, dot, fill=C_TEAL)
            self._rich_textbox(slide, right_x + 0.52, y,
                               right_w - 0.72, item_h_right,
                               it, size=13, color=body_color, font=F_BODY)
        # Principle (bottom of right panel)
        if principle:
            p_y = ct + ch - principle_h + 0.02
            self._hline(slide, right_x + 0.24, p_y, right_w - 0.48,
                        C_MGRAY if not dark else C_TEAL, 0.02)
            tb = slide.shapes.add_textbox(
                I(right_x + 0.24), I(p_y + 0.10),
                I(right_w - 0.48), I(principle_h - 0.14))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = principle
            run.font.name = F_BODY
            run.font.size = Pt(12)
            run.font.italic = True
            run.font.color.rgb = C_LGRAY if dark else C_GRAY
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def not_list(self, title: str, items: list, dark: bool = False,
                 eyebrow: str = '', subtitle: str = '',
                 footnote: str = '', notes: str = ''):
        """Grid of 'what is / proper role' rows.

        items: list of dicts {what, correct}. 2-7 rows."""
        n = len(items)
        assert 2 <= n <= 7, 'not_list requires 2-7 items'
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        row_h = ch / n
        for i, item in enumerate(items):
            y = ct + i * row_h
            # Index (big gray mono)
            self._textbox(slide, CL, y + 0.05, 1.00, row_h - 0.10,
                          f'{i+1:02d}', size=28, color=C_GRAY,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # what — 15pt navy bold
            what_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, CL + 1.05, y + 0.08, 5.50, 0.45,
                          item.get('what', ''), size=15, bold=True,
                          color=what_color, font=F_BODY, align=PP_ALIGN.LEFT)
            # PROPER ROLE label
            self._textbox(slide, CL + 6.65, y + 0.08, 1.60, 0.32,
                          'PROPER ROLE', size=9, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # correct body
            body_color = C_LGRAY if dark else C_GRAY
            self._rich_textbox(slide, CL + 6.65, y + 0.40, CW - 6.70, row_h - 0.48,
                               item.get('correct', ''), size=13,
                               color=body_color, font=F_BODY)
            # Separator line between rows (not below last)
            if i < n - 1:
                self._hline(slide, CL, y + row_h - 0.02, CW,
                            C_MGRAY if not dark else C_TEAL, 0.015)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def before_after(self, title: str, before: dict, after: dict,
                     dark: bool = False,
                     eyebrow: str = '', subtitle: str = '',
                     footnote: str = '', notes: str = ''):
        """FROM · 过去 → TO · 未来 three-column layout.

        before / after: dicts {title, body}."""
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        left_w  = CW * 0.44
        arrow_w = 0.60
        right_w = CW - left_w - arrow_w
        mid_x   = CL + left_w
        right_x = mid_x + arrow_w

        muted_title = C_GRAY if not dark else C_LGRAY
        muted_body  = C_GRAY if not dark else C_LGRAY

        # FROM (left)
        self._textbox(slide, CL, ct + 0.10, left_w, 0.28,
                      'FROM · 过去', size=11, bold=True, color=C_GRAY,
                      font=F_ACCENT, align=PP_ALIGN.LEFT)
        self._textbox(slide, CL, ct + 0.42, left_w, 0.80,
                      before.get('title', ''), size=22, bold=True,
                      color=muted_title, font=F_BODY, align=PP_ALIGN.LEFT)
        self._rich_textbox(slide, CL, ct + 1.38, left_w, ch - 1.50,
                           before.get('body', ''),
                           size=14, color=muted_body, font=F_BODY)

        # Arrow (center)
        self._textbox(slide, mid_x, ct + ch / 2 - 0.40, arrow_w, 0.80,
                      '→', size=36, color=C_TEAL,
                      font=F_ACCENT, align=PP_ALIGN.CENTER)

        # TO (right) — light card + thick navy left border
        self._rect(slide, right_x, ct, right_w, ch, fill=C_LGRAY, radius=0.06)
        self._rect(slide, right_x, ct, 0.08, ch, fill=C_NAVY)
        self._textbox(slide, right_x + 0.24, ct + 0.10, right_w - 0.36, 0.28,
                      'TO · 未来', size=11, bold=True, color=C_TEAL,
                      font=F_ACCENT, align=PP_ALIGN.LEFT)
        self._textbox(slide, right_x + 0.24, ct + 0.42, right_w - 0.36, 0.80,
                      after.get('title', ''), size=22, bold=True,
                      color=C_NAVY, font=F_BODY, align=PP_ALIGN.LEFT)
        self._rich_textbox(slide, right_x + 0.24, ct + 1.38, right_w - 0.36, ch - 1.50,
                           after.get('body', ''),
                           size=14, color=C_BLACK, font=F_BODY)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def value_ladder(self, title: str, stages: list, caption: str = '',
                     dark: bool = False,
                     eyebrow: str = '', subtitle: str = '',
                     footnote: str = '', notes: str = ''):
        """3-4 stage columns with STAGE label, title, tag, progress bar, state.

        stages: list of dicts {title, tag, progress (0-100), state}."""
        n = len(stages)
        assert 3 <= n <= 4, 'value_ladder requires 3-4 stages'
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        chevron_w = 0.50
        total_gaps = chevron_w * (n - 1)
        col_w = (CW - total_gaps) / n

        for i, stage in enumerate(stages):
            x = CL + i * (col_w + chevron_w)
            # STAGE label
            self._textbox(slide, x, ct + 0.05, col_w, 0.28,
                          f'STAGE {i+1:02d}', size=11, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # Stage title
            title_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, x, ct + 0.32, col_w, 0.55,
                          stage.get('title', ''), size=16, bold=True,
                          color=title_color, font=F_BODY, align=PP_ALIGN.LEFT)
            # Tag
            if stage.get('tag'):
                tag_color = C_LGRAY if dark else C_GRAY
                self._textbox(slide, x, ct + 0.90, col_w, 0.35,
                              stage.get('tag', ''), size=12, color=tag_color,
                              font=F_BODY, align=PP_ALIGN.LEFT)
            # Progress bar
            prog = max(0, min(100, int(stage.get('progress', 0))))
            pb_y = ct + 1.38
            pb_h = 0.10
            self._rect(slide, x, pb_y, col_w, pb_h, fill=C_MGRAY)
            fill_w = col_w * prog / 100
            fill_color = C_NAVY if prog == 100 else C_TEAL
            if fill_w > 0:
                self._rect(slide, x, pb_y, fill_w, pb_h, fill=fill_color)
            # State label
            state_color = C_LGRAY if dark else C_GRAY
            self._textbox(slide, x, pb_y + 0.20, col_w, 0.35,
                          stage.get('state', ''), size=12, color=state_color,
                          font=F_BODY, align=PP_ALIGN.LEFT)
            # Chevron between columns
            if i < n - 1:
                cx = x + col_w
                self._textbox(slide, cx, ct + 0.30, chevron_w, 0.80,
                              '›', size=32, color=C_GRAY,
                              font=F_ACCENT, align=PP_ALIGN.CENTER)
        # Legacy caption (bottom-center) only when no footnote
        if caption and not footnote:
            self._textbox(slide, CL, ct + ch - 0.30, CW, 0.28,
                          caption, size=12, color=C_GRAY,
                          font=F_BODY, align=PP_ALIGN.CENTER)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def big_number(self, title: str, number: str, unit: str = '',
                   label: str = '', body: str = '',
                   dark: bool = False,
                   eyebrow: str = '', subtitle: str = '',
                   footnote: str = '', notes: str = ''):
        """Left: giant navy number + small teal unit. Right: label + body."""
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        n_len = len(str(number))
        if   n_len <= 2: num_size = 150
        elif n_len <= 4: num_size = 120
        elif n_len <= 6: num_size = 90
        elif n_len <= 8: num_size = 70
        else:            num_size = 50

        left_w  = CW * 0.50
        right_x = CL + left_w + 0.20
        right_w = CW - left_w - 0.20

        num_color = C_WHITE if dark else C_NAVY
        self._textbox(slide, CL, ct + 0.20, left_w, ch - 0.30,
                      str(number), size=num_size, bold=True,
                      color=num_color, font=F_ACCENT,
                      align=PP_ALIGN.LEFT, wrap=False)
        if unit:
            self._textbox(slide, CL + left_w - 1.20, ct + 0.80, 1.20, 0.50,
                          unit, size=20, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
        # Right: short teal rule + label + body
        self._rect(slide, right_x, ct + 0.30, 0.80, 0.05, fill=C_TEAL)
        if label:
            self._textbox(slide, right_x, ct + 0.42, right_w, 0.32,
                          label.upper(), size=11, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
        body_color = C_WHITE if dark else C_NAVY
        if body:
            self._rich_textbox(slide, right_x, ct + 0.82, right_w, ch - 0.92,
                               body, size=19, color=body_color, font=F_BODY)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def quote(self, text: str, attribution: str = '', title: str = '',
              dark: bool = True,
              eyebrow: str = '', subtitle: str = '',
              footnote: str = '', notes: str = ''):
        """Large pull-quote with giant teal ', optional title on top."""
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        if title:
            self._title(slide, title, dark=dark, eyebrow=eyebrow)
            self._subtitle_block(slide, subtitle, dark)
            ct, _, ch = self._content_bounds(subtitle, footnote, None)
            y0 = ct + 0.05
        elif eyebrow:
            self._textbox(slide, TL, EYEBROW_Y, TITLE_SAFE_W, EYEBROW_H,
                          eyebrow.upper(), size=11, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            y0 = 0.90
        else:
            y0 = 0.90
        # Giant teal opening quote mark
        self._textbox(slide, CL, y0, 2.00, 1.80,
                      '"', size=100, bold=True, color=C_TEAL,
                      font=F_ACCENT, align=PP_ALIGN.LEFT)
        # Quote body
        body_color = C_WHITE if dark else C_NAVY
        body_y = y0 + 1.20
        body_h = 3.50
        self._rich_textbox(slide, CL + 0.40, body_y, CW - 0.50, body_h,
                           text, size=26, color=body_color, font=F_BODY)
        # Attribution
        if attribution:
            self._textbox(slide, CL + 0.40, body_y + body_h + 0.10, CW - 0.50, 0.35,
                          f'— {attribution}', size=12, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def number_list(self, title: str, items: list, dark: bool = False,
                    eyebrow: str = '', subtitle: str = '',
                    footnote: str = '', notes: str = ''):
        """Rows with big teal mono numeral + title + body.

        items: list of dicts {title, body}. 3-5 rows."""
        n = len(items)
        assert 3 <= n <= 5, 'number_list requires 3-5 items'
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        row_h = ch / n
        for i, item in enumerate(items):
            y = ct + i * row_h
            # Big numeral
            self._textbox(slide, CL, y + 0.05, 1.20, row_h - 0.10,
                          f'{i+1:02d}', size=36, bold=True, color=C_TEAL,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # Item title
            title_color = C_WHITE if dark else C_NAVY
            self._textbox(slide, CL + 1.30, y + 0.10, CW - 1.40, 0.48,
                          item.get('title', ''), size=17, bold=True,
                          color=title_color, font=F_BODY, align=PP_ALIGN.LEFT)
            # Body
            body_color = C_LGRAY if dark else C_GRAY
            self._rich_textbox(slide, CL + 1.30, y + 0.58, CW - 1.40, row_h - 0.66,
                               item.get('body', ''), size=14,
                               color=body_color, font=F_BODY)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def step_grid(self, title: str, steps: list,
                  highlight_last: bool = True, dark: bool = False,
                  eyebrow: str = '', subtitle: str = '',
                  footnote: str = '', notes: str = ''):
        """3-6 card grid with step number + title + body.

        steps: list of dicts {title, body}. If highlight_last and not dark,
        last card gets C_NAVY bg + white/lgray text."""
        n = len(steps)
        assert 3 <= n <= 6, 'step_grid requires 3-6 steps'
        slide = self._new_slide('custom slide1-light', dark=dark, density='compact')
        self._title(slide, title, dark=dark, eyebrow=eyebrow)
        self._subtitle_block(slide, subtitle, dark)

        ct, _, ch = self._content_bounds(subtitle, footnote, None)

        # Grid: prefer single row for 3, 2 rows for 4+
        if n <= 3:
            cols = n; rows = 1
        elif n == 4:
            cols = 4; rows = 1
        elif n == 5:
            cols = 5; rows = 1
        else:  # 6
            cols = 3; rows = 2

        GAP = 0.20
        card_w = (CW - GAP * (cols - 1)) / cols
        card_h = max(2.5, (ch - GAP * (rows - 1)) / rows)

        for i, step in enumerate(steps):
            col_i = i % cols
            row_i = i // cols
            x = CL + col_i * (card_w + GAP)
            y = ct + row_i * (card_h + GAP)

            is_highlighted = (highlight_last and not dark and i == n - 1)
            card_bg = C_NAVY if is_highlighted else (
                RGBColor(0x1A, 0x2E, 0x7A) if dark else C_LGRAY)
            accent = C_TEAL if not is_highlighted else RGBColor(0x5B, 0xF2, 0xF5)

            # Card bg (rounded)
            self._rect(slide, x, y, card_w, card_h, fill=card_bg, radius=0.10)
            # Top teal accent strip
            self._rect(slide, x, y, card_w, 0.06, fill=accent)
            # Step number
            num_color = accent
            self._textbox(slide, x + 0.22, y + 0.16, card_w - 0.30, 0.36,
                          f'{i+1:02d}', size=16, bold=True, color=num_color,
                          font=F_ACCENT, align=PP_ALIGN.LEFT)
            # Card title
            title_color = C_WHITE if (is_highlighted or dark) else C_NAVY
            self._textbox(slide, x + 0.22, y + 0.55, card_w - 0.30, 0.60,
                          step.get('title', ''), size=20, bold=True,
                          color=title_color, font=F_BODY, align=PP_ALIGN.LEFT)
            # Card body
            body_color = C_LGRAY if (is_highlighted or dark) else C_GRAY
            self._rich_textbox(slide, x + 0.22, y + 1.18, card_w - 0.30,
                               card_h - 1.30, step.get('body', ''),
                               size=14, color=body_color, font=F_BODY)
        self._footnote_strip(slide, footnote, dark)
        self._attach_notes(slide, notes)
        return slide

    def save(self) -> str:
        os.makedirs(os.path.dirname(os.path.abspath(self.output)),
                    exist_ok=True)
        self.prs.save(self.output)
        return self.output
